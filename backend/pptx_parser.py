"""
PPTX Parser Module

Extracts structured slide data (titles, bullets) from PowerPoint files.
This is different from file_utils.py which extracts raw text.
"""

import logging
from typing import List, Dict
from pptx import Presentation

logger = logging.getLogger(__name__)


def extract_slides_structured(pptx_path: str) -> List[Dict[str, str]]:
    """
    Extract structured slide data from PPTX file
    
    Args:
        pptx_path: Path to .pptx file
        
    Returns:
        List of slides with 'title' and 'bullets' keys
        
    Example:
        [
            {
                "title": "Introduction to Cloud",
                "bullets": "• AWS services\n• Azure platform\n• Google Cloud"
            },
            ...
        ]
    """
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            title = ""
            bullets = []
            
            # Extract title
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    
                    # Skip if this is the title shape
                    if shape == slide.shapes.title:
                        continue
                    
                    # Extract paragraphs
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            bullets.append(text)
            
            # Only include slides with content
            if title or bullets:
                slide_data = {
                    "title": title or f"Slide {slide_num}",
                    "bullets": "\n".join(bullets)
                }
                slides_data.append(slide_data)
                logger.debug(f"Slide {slide_num}: {title} ({len(bullets)} bullets)")
        
        logger.info(f"Extracted {len(slides_data)} slides from PPTX")
        return slides_data
        
    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")
        raise RuntimeError(f"PPTX parsing failed: {str(e)}")


def extract_slides_simple(pptx_path: str) -> List[Dict[str, str]]:
    """
    Simplified extraction - just get titles and combine all text
    
    Args:
        pptx_path: Path to .pptx file
        
    Returns:
        List of slides with 'title' and 'bullets' keys
    """
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            title = ""
            all_text = []
            
            # Get all text from slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # First non-empty text is usually the title
                        if not title and shape == slide.shapes.title:
                            title = text
                        else:
                            all_text.append(text)
            
            # Only include slides with content
            if title or all_text:
                slide_data = {
                    "title": title or f"Slide {slide_num}",
                    "bullets": "\n".join(all_text)
                }
                slides_data.append(slide_data)
        
        logger.info(f"Extracted {len(slides_data)} slides (simple mode)")
        return slides_data
        
    except Exception as e:
        logger.error(f"Simple PPTX parsing failed: {e}")
        raise RuntimeError(f"PPTX parsing failed: {str(e)}")