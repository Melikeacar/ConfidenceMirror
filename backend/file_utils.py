from pptx import Presentation
from pypdf import PdfReader
import logging
import os

logger = logging.getLogger(__name__)

def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from a PowerPoint file
    
    Args:
        file_path: Path to .pptx file
        
    Returns:
        Extracted text content
    """
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            
            # Add slide header
            if slide_text:
                text_content.append(f"--- Slide {i+1} ---")
                text_content.extend(slide_text)
                
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to extract PPTX text: {e}")
        raise ValueError("Could not parse PowerPoint file")

def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text from a PDF file
    
    Args:
        file_path: Path to .pdf file
        
    Returns:
        Extracted text content
    """
    try:
        reader = PdfReader(file_path)
        text_content = []
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                text_content.append(f"--- Page {i+1} ---")
                text_content.append(text)
                
        return "\n".join(text_content)
    except Exception as e:
        logger.error(f"Failed to extract PDF text: {e}")
        raise ValueError("Could not parse PDF file")
